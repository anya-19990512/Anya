from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
maroon = RGBColor(139, 30, 45)
teal = RGBColor(15, 118, 110)
blue = RGBColor(30, 58, 138)
ink = RGBColor(17, 24, 39)
muted = RGBColor(91, 100, 114)
light = RGBColor(247, 248, 250)


def add_title_slide(title, subtitle1, subtitle2=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(247, 248, 250)

    # top bar
    left = Inches(0.0)
    top = Inches(0.0)
    width = prs.slide_width
    height = Inches(0.35)
    shapes = slide.shapes
    bar = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = maroon
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(2.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ink
    p.alignment = PP_ALIGN.LEFT

    if subtitle1:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(9.2), Inches(1.2))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle1
        p2.font.size = Pt(16)
        p2.font.color.rgb = muted
        p2.alignment = PP_ALIGN.LEFT
    if subtitle2:
        sub_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(9.2), Inches(1.0))
        tf3 = sub_box2.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = subtitle2
        p3.font.size = Pt(16)
        p3.font.color.rgb = muted
        p3.alignment = PP_ALIGN.LEFT

    # accent panel
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(1.4), Inches(3.2), Inches(3.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(217, 222, 231)
    panel.line.width = Pt(1.2)

    panel_tf = panel.text_frame
    panel_tf.word_wrap = True
    p4 = panel_tf.paragraphs[0]
    p4.text = '跨領域學院\n合作與創新入口'
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = maroon
    p4.alignment = PP_ALIGN.CENTER


def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = light

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), prs.slide_width, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = teal
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(9.6), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ink

    body = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(10.6), Inches(4.8))
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.bullet = True
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT

    # side panel
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(1.4), Inches(1.8), Inches(3.6))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(217, 222, 231)
    panel.line.width = Pt(1.0)
    panel_tf = panel.text_frame
    p = panel_tf.paragraphs[0]
    p.text = '合作\n入口'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = blue
    p.alignment = PP_ALIGN.CENTER


add_title_slide('台北醫學大學跨領域學院', '對外合作與學生學習資源簡報', '面向醫療、科技、創新與跨域實作')

slides = [
    ('醫療創新已經不只是醫療端的問題', [
        'AI、資料、永續、高齡照護與社區健康都需要跨學科合作。',
        '跨領域學院作為校內外合作的共同入口。'
    ]),
    ('我們的定位', [
        '將醫療與健康議題轉譯為可教、可學、可合作的跨域方案。',
        '讓學生與教師能以真實場域與資源進行實作。'
    ]),
    ('學院成立背景', [
        '2018 年成立第十一個學院，專司發展醫療外知識與技能領域。',
        '作為專業與臨床教育之間的溝通橋樑。'
    ]),
    ('我們不是單一課程單位，而是跨域合作平台', [
        '把校內外資源轉譯為學生可使用、教師可共創、外部單位可合作的系統。',
        '讓跨域教育不只停留在課程，而是形成合作機制。'
    ]),
    ('我們培養能與外部單位共同解題的人', [
        '主動探索、創新實踐、團隊整合、全球適應與自主學習。',
        '培養能面對複雜議題的跨域能力。'
    ]),
    ('跨領域學習中心', [
        '規劃跨領域課程、微學程與教師社群。',
        '建立系統化學習與跨校合作路徑。'
    ]),
    ('創新創業教育中心', [
        '以課程與輔導培育創新能力。',
        '銜接校內外新創資源與團隊進駐。'
    ]),
    ('數位自學中心', [
        '導入國際優質磨課師課程。',
        '支援自主學習與全球共學。'
    ]),
    ('智慧醫療跨領域學士學位學程', [
        '以人工智慧為導向，結合生醫知識與實際應用。',
        '提供面向 AI Healthcare 的完整學習路徑。'
    ]),
    ('學院資源可被使用', [
        '杏春樓 1F / B1F、VR LAB、設備借用與團隊進駐。',
        '讓資源從展示轉為合作與實作工具。'
    ]),
    ('近期亮點', [
        '社會處方箋 MOU、暑期數位自學遊學營與 5G 加速器活動。',
        '提供外部單位多元切入合作的機會。'
    ]),
    ('國際與開放教育合作', [
        '連結九州大學、Stanford Biodesign 與開放教育資源。',
        '強化學生跨國團隊協作與全球議題能力。'
    ]),
    ('外部單位可以從四種合作開始', [
        '議題合作、課程合作、活動合作與實作場域合作。',
        '降低合作門檻，讓與學院接軌更容易。'
    ]),
    ('結語', [
        '跨領域學院歡迎外部單位以不同方式加入合作。',
        '讓醫療創新合作，從一個清楚的入口開始。'
    ]),
]

for title, bullets in slides:
    add_content_slide(title, bullets)

out_path = 'tmu_interdisciplinary_college_deck.pptx'
prs.save(out_path)
print(out_path)
